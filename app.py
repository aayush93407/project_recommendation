from flask import Flask, request, render_template, jsonify
from flask_cors import CORS
import os
from pdfminer.high_level import extract_text
from pptx import Presentation
from mistralai import Mistral

# 🔥 Initialize Flask App
app = Flask(__name__)
CORS(app)  # Enable Cross-Origin Requests

# 🔥 API Key (HARD-CODED)
API_KEY_MISTRAL = "aKFEMuDwJOvtphHDDOrh2qbfRP7jEA1L"

# 🔥 Ensure 'uploads' directory exists
UPLOAD_FOLDER = "uploads"
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# ✅ Extract text from PDF
def extract_text_pdf(pdf_path):
    try:
        return extract_text(pdf_path)
    except Exception as e:
        print(f"Error extracting PDF: {e}")
        return ""

# ✅ Extract text from PPT
def extract_text_ppt(ppt_path):
    try:
        ppt = Presentation(ppt_path)
        text = "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
        return text
    except Exception as e:
        print(f"Error extracting PPT: {e}")
        return ""

# ✅ Analyze Resume with Mistral AI
def analyze_resume(text):
    prompt = f"""
    Extract the following details from the resume text:
    - Name
    - Email ID
    - Contact Number
    - LinkedIn
    - GitHub
    - Skills (comma-separated)
    Format as:
    Name: <value>
    Email ID: <value>
    Contact Number: <value>
    LinkedIn: <value>
    GitHub: <value>
    Skills: <comma-separated list>
    Resume Text:
    {text}
    """
    try:
        client = Mistral(api_key=API_KEY_MISTRAL)
        response = client.chat.complete(model="mistral-large-latest", messages=[{"role": "user", "content": prompt}])
        return response.choices[0].message.content
    except Exception as e:
        print(f"Error with Mistral API: {e}")
        return "Error processing resume."

# ✅ Recommend Projects Based on Skills
def recommend_projects(skills):
    if not skills:
        return "No skills detected, unable to suggest projects."
    
    prompt = f"Suggest 5 projects based on these skills: {skills}."
    try:
        client = Mistral(api_key=API_KEY_MISTRAL)
        response = client.chat.complete(model="mistral-large-latest", messages=[{"role": "user", "content": prompt}])
        return response.choices[0].message.content
    except Exception as e:
        print(f"Error with Mistral API: {e}")
        return "Error generating project recommendations."

# ✅ File Upload Route
@app.route("/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    if file.filename == "":
        return jsonify({"error": "No selected file"}), 400

    file_ext = os.path.splitext(file.filename)[1].lower()
    file_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(file_path)

    # 🔥 Extract text from the uploaded file
    if file_ext == ".pdf":
        resume_text = extract_text_pdf(file_path)
    elif file_ext in [".ppt", ".pptx"]:
        resume_text = extract_text_ppt(file_path)
    else:
        return jsonify({"error": "Invalid file format"}), 400

    if not resume_text.strip():
        return jsonify({"error": "Failed to extract text from file"}), 400

    # 🔥 Analyze Resume and Get Information
    resume_info = analyze_resume(resume_text)
    resume_data = {line.split(":")[0].strip(): line.split(":")[1].strip() for line in resume_info.split("\n") if ":" in line}

    # 🔥 Get Project Recommendations
    skills = resume_data.get("Skills", "")
    resume_data["Projects"] = recommend_projects(skills)

    return jsonify(resume_data)

# ✅ Serve the Upload Page
@app.route("/")
def index():
    return render_template("index.html")

# ✅ Serve the Results Page
@app.route("/result")
def result():
    return render_template("result.html")

if __name__ == "__main__":
    app.run(debug=True)
